from langgraph.graph import START, END, StateGraph
from pydantic import BaseModel, Field
from langchain_google_genai import ChatGoogleGenerativeAI
from typing import Literal, TypedDict, List, Annotated, Optional
from langchain_core.messages import SystemMessage, HumanMessage
import os
import operator
from langgraph.types import Send
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from langchain_groq import ChatGroq
import io
from dotenv import load_dotenv
from langchain_tavily import TavilySearch
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from huggingface_hub import InferenceClient
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import time

load_dotenv()

gemini_model = ChatGoogleGenerativeAI(model="gemini-2.5-flash", google_api_key=os.getenv("GOOGLE_API_KEY_PPT_AGENT"))

grok_modal = ChatGroq(
    model="llama-3.3-70b-versatile",
    api_key=os.getenv("GROQ_API_KEY"),
    temperature=0.1, # Keep temperature low for strict JSON formatting
)


class SlideTask(BaseModel):
    id: int
    title: str
    goal: str = Field(..., description="What the audience should learn from this slide.")
    requires_research: bool = False

class PresentationPlan(BaseModel):
    ppt_title: str
    audience: str = Field(..., description="Target audience for this presentation.")
    tone: str = Field(..., description="Speaking tone (e.g., persuasive, educational).")
    slides: List[SlideTask] = Field(..., max_length=10, description="MAXIMUM 10 SLIDES.")


class SlideContent(BaseModel):
    slide_id: int
    title: str
    bullet_points: List[str] = Field(..., min_length=3, max_length=5, description="3 to 5 short, punchy bullet points.")
    speaker_notes: str = Field(..., description="Detailed paragraph for the presenter to read.")


class EvidenceItem(BaseModel):
    title: str
    url: str
    snippet: Optional[str] = None


class EvidencePack(BaseModel):
    evidence: List[EvidenceItem] = Field(default_factory=list)


class RouterDecision(BaseModel):
    needs_research: bool
    queries: List[str] = Field(default_factory=list)

class ImageSpec(BaseModel):
    slide_id: int = Field(..., description="The ID of the slide where this image should go.")
    filename: str = Field(..., description="Save under images/, e.g. flow_diagram.png")
    prompt: str = Field(..., description="Highly detailed prompt to send to the image generation model. No text.")

class GlobalImagePlan(BaseModel):
    images: List[ImageSpec] = Field(
        default_factory=list, 
        max_length=4, # <--- STRICTLY LIMIT TO MAX 4 IMAGES OVERALL
        description="Select up to 4 slides that need visual representation."
    )

class State(TypedDict):
    topic: str
    saved_filepath:str
    needs_research: bool
    queries: List[str]
    evidence: List[EvidenceItem]
    plan: Optional[PresentationPlan]
    slides_content: Annotated[List[SlideContent], operator.add] 
    image_specs: List[dict]

ROUTER_SYSTEM = """You are the routing module for an Enterprise Presentation Generator.
Decide if web research is required to create highly accurate slides. Set needs_research=True if facts/news are needed."""

RESEARCH_SYSTEM = """You are a Research Synthesizer. Extract the most relevant facts and URLs from the raw search results to create an EvidencePack."""

ORCH_SYSTEM = """You are a Lead Presentation Architect. Produce a comprehensive outline for a premium PowerPoint presentation.
Requirements:
1. Create a maximum of 10 slides.
2. Ensure a logical flow (Introduction -> Core Concepts -> Analysis -> Conclusion).
Output STRICTLY matching the PresentationPlan schema."""

WORKER_SYSTEM = """You are an Expert Presentation Designer. Write the content for ONE slide.
Constraints:
- Provide 3 to 5 short, punchy bullet points. DO NOT write long paragraphs for the bullets.
- Write highly detailed speaker notes explaining the concepts deeply.
- If evidence is provided, weave the facts naturally into the bullets and notes."""

DECIDE_IMAGES_SYSTEM = """You are an expert visual editor. Decide which slides need images.
CRITICAL RULES:
1. MAXIMUM 4 IMAGES TOTAL for the entire presentation.
2. Map the image to the exact slide_id.
3. NEVER ask for "text", "infographics", or "labels". The AI cannot spell.
4. Ask for "Clean abstract 3D render" or "Conceptual vector illustration"."""

def router_node(state: State):
    router_decision = gemini_model.with_structured_output(RouterDecision)
    result = router_decision.invoke([
        SystemMessage(content=ROUTER_SYSTEM),
        HumanMessage(content=f"Topic: {state['topic']}")
    ])
    return {"needs_research": result.needs_research, "queries": result.queries}

def route(state:State)->str:
    if state.get("needs_research",False):
        return "research"
    else:
        return "orchestrator"
    
def research(state:State)->dict:
    print("--- RESEARCHING ---")
    tool=TavilySearch(max_results=3)
    web_search_result=[]
    for query in state.get("queries",[]):
        web_search_result.extend(tool.invoke(query))
    
    if not web_search_result:
        return {"evidence": []}

    evidencepack=gemini_model.with_structured_output(EvidencePack)
    result = evidencepack.invoke([
        SystemMessage(content=RESEARCH_SYSTEM),
        HumanMessage(content=f"Raw Results : {web_search_result}")
    ])
    dedup = {e.url: e for e in result.evidence if e.url}
    return {"evidence": list(dedup.values())}

def orchestrator(state: State) -> dict:
    print("--- PLANNING PRESENTATION ---")
    planner = gemini_model.with_structured_output(PresentationPlan)
    answer = planner.invoke([
        SystemMessage(content=ORCH_SYSTEM),
        HumanMessage(content=f"Topic: {state['topic']}")
    ])
    return {"plan": answer}

def n_worker_excution(state:State):
    tasks=[]
    for task in state["plan"].slides:
        tasks.append(Send("worker",{"task":task.model_dump(),"evidence":[i.model_dump() for i in state.get("evidence", [])],"plan_context":state["plan"].model_dump()}))
    return tasks

def worker(payload:dict)->dict:
    task=SlideTask(**payload["task"])
    evidence=[EvidenceItem(**e) for e in payload.get("evidence", [])]
    plan=PresentationPlan(**payload["plan_context"])

    print(f"    -> Designing Slide: {task.title}")

    if evidence:
        evidence_text = "\n".join([f"- {e.title}: {e.url}" for e in evidence]) 
    else:
         evidence_text="None"

    print("--- Worker node ---")
    slide_writer = grok_modal.with_structured_output(SlideContent)
    slide_content = slide_writer.invoke([
        SystemMessage(content=WORKER_SYSTEM),
        HumanMessage(content=
        f"""
        Slide ID: {task.id}
        Slide Title: {task.title}
        Goal: {task.goal}
        Available Evidence: {evidence_text}
        Generate the bullet points and speaker notes for this slide.
        """)
    ])

    return {"slides_content": [slide_content]}

def decide_images(state: State) -> dict:
    print("--- AI EDITOR: PLANNING IMAGES ---")
    
    # Create a summary of the slides so the image editor knows what they are about
    slide_summaries = "\n".join([f"Slide {s.slide_id}: {s.title}" for s in state['slides_content']])
    
    planner = gemini_model.with_structured_output(GlobalImagePlan)
    image_plan = planner.invoke([
        SystemMessage(content=DECIDE_IMAGES_SYSTEM),
        HumanMessage(content=f"Topic: {state['topic']}\n\nSlides Overview:\n{slide_summaries}")
    ])
    
    specs = [img.model_dump() for img in image_plan.images][:4] 
    return {"image_specs": specs}

# hf_token = os.environ.get("HUGGINGFACE_IMAGE_GENERATION_PPT_AGENT_API_TOKEN_SANJANA")
# hf_token

def _generate_image_bytes(prompt: str) -> bytes:
    print(f"      -> Painting (SDXL): '{prompt}'")
    hf_token = os.environ.get("HUGGINGFACE_IMAGE_GENERATION_PPT_AGENT_API_TOKEN_SANJANA")
    # client = InferenceClient(provider="hf-inference", api_key=hf_token)
    # image = client.text_to_image(prompt, model="stabilityai/stable-diffusion-xl-base-1.0")
    try:
        client = InferenceClient(
        model="stabilityai/stable-diffusion-xl-base-1.0",
        token=hf_token
        )
       
        image = client.text_to_image(prompt)
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='PNG')
        return img_byte_arr.getvalue()
    except Exception as e:
        print(f"      -> ❌ Image failed: {e}")
        return None


def generate_and_build_pptx(state: State) -> dict:
    print("--- BUILDING GAMMA-STYLE POWERPOINT ---")
    plan = state["plan"]
    slides_content = sorted(state["slides_content"], key=lambda x: x.slide_id)
    image_specs = state.get("image_specs", [])
    
    images_dir = Path("images")
    images_dir.mkdir(exist_ok=True)

    # 1. Generate Images
    slide_images = {}
    for spec in image_specs:
        filename = "".join(c if c.isalnum() or c == "." else "_" for c in spec["filename"])
        out_path = images_dir / filename
        if not out_path.exists():
            img_bytes = _generate_image_bytes(spec["prompt"])
            if img_bytes:
                out_path.write_bytes(img_bytes)
        if out_path.exists():
            slide_images[spec['slide_id']] = str(out_path)

    # 2. Build PPTX
    prs = Presentation()
    
    # Premium Color Palette (Gamma Dark Mode)
    BG_COLOR = RGBColor(18, 18, 20)         # Very Dark Grey Background
    CARD_COLOR = RGBColor(30, 30, 34)       # Slightly lighter grey for text blocks
    ACCENT_COLOR = RGBColor(99, 102, 241)   # Vibrant Indigo/Purple Accent
    TITLE_COLOR = RGBColor(255, 255, 255)   # Pure White
    SUB_COLOR = RGBColor(161, 161, 170)     # Muted Grey for subtext
    TEXT_COLOR = RGBColor(228, 228, 231)    # Off-white for readability

    # --- Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # Decorative Accent Bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.5), Inches(0.1), Inches(2.0))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_COLOR
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(8.0), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = plan.ppt_title
    p.font.size = Pt(44) if len(plan.ppt_title) < 40 else Pt(36) # Dynamic Title Size
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Subtitle / Details
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(8.0), Inches(1.0))
    sub_p = sub_box.text_frame.add_paragraph()
    sub_p.text = "Generated by AI Agent"
    sub_p.font.size = Pt(20)
    sub_p.font.color.rgb = SUB_COLOR

    # --- Content Slides ---
    for sc in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout for total control
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG_COLOR
        
        has_image = sc.slide_id in slide_images
        
        # 1. Dynamic Title Sizing
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(1.0))
        title_p = title_box.text_frame.add_paragraph()
        title_p.text = sc.title
        title_p.font.bold = True
        title_p.font.color.rgb = TITLE_COLOR
        
        # Adjust size based on character count to prevent overflow
        if len(sc.title) > 45:
            title_p.font.size = Pt(28)
        elif len(sc.title) > 30:
            title_p.font.size = Pt(32)
        else:
            title_p.font.size = Pt(40)

        # 2. Add an accent line under the title
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(2.0), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_COLOR
        line.line.fill.background()

        # 3. Calculate Body Text Size dynamically based on content length
        total_chars = sum([len(b) for b in sc.bullet_points])
        if total_chars > 400:
            bullet_size = Pt(16)
            spacing = Pt(12)
        elif total_chars > 250:
            bullet_size = Pt(18)
            spacing = Pt(16)
        else:
            bullet_size = Pt(22)
            spacing = Pt(20)

        # 4. Text Layout (With or Without Image)
        text_width = Inches(4.2) if has_image else Inches(8.5)
        
        # Create a "Card" background for the text (Gamma style)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), text_width, Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.fill.background() # Remove border

        # Create Text Box over the card
        body_shape = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), text_width - Inches(0.2), Inches(4.3))
        tf = body_shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE # Center text vertically in the box
        
        # 5. Add Custom Symbols to Bullet Points
        for bullet in sc.bullet_points:
            p = tf.add_paragraph()
            p.text = f"✦  {bullet}" # Custom modern symbol
            p.font.color.rgb = TEXT_COLOR
            p.font.size = bullet_size
            p.space_after = spacing
            p.level = 0

        # Add Speaker Notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = sc.speaker_notes

        # 6. Perfect Image Alignment
        if has_image:
            img_path = slide_images[sc.slide_id]
            slide.shapes.add_picture(
                img_path, 
                left=Inches(5.0),   
                top=Inches(1.8),    
                width=Inches(4.5),
                height=Inches(4.5)  # Forces a perfect square alignment
            )

    # Save
    timestamp = int(time.time())
    safe_title = "".join(c if c.isalnum() else "_" for c in plan.ppt_title).strip("_").lower()
    filename = f"{safe_title}_{timestamp}.pptx"
    export_dir = os.path.join(os.getcwd(), "exports")
    os.makedirs(export_dir, exist_ok=True)
    
    # Path for the Hard Drive
    save_path = os.path.join(export_dir, filename)
    
    # Save to the specific folder
    prs.save(save_path)
    print(f"\n✅ SUCCESS! Gamma-Style PowerPoint saved as: {save_path}")

    # Return the web-friendly relative path for the UI button
    web_path = f"exports/{filename}"

    return {"saved_filepath": web_path}

graph = StateGraph(State)
graph.add_node("orchestrator", orchestrator)
graph.add_node("research", research)
graph.add_node("router_node", router_node)
graph.add_node("worker", worker)
graph.add_node("decide_images", decide_images)
graph.add_node("generate_and_build_pptx", generate_and_build_pptx)

graph.add_edge(START, "router_node")
graph.add_conditional_edges("router_node", route, {"orchestrator": "orchestrator", "research": "research"})
graph.add_edge("research", "orchestrator")
graph.add_conditional_edges("orchestrator", n_worker_excution, ["worker"])
graph.add_edge("worker", "decide_images")
graph.add_edge("decide_images", "generate_and_build_pptx")
graph.add_edge("generate_and_build_pptx", END)

ppt_workflow=graph.compile()












