# utils/multimodal_processor.py
import os
import io
import uuid
import base64
import hashlib
import zipfile
from PIL import Image
from pptx import Presentation
from unstructured.partition.auto import partition
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.messages import HumanMessage
from langchain_groq import ChatGroq
from langchain_text_splitters import RecursiveCharacterTextSplitter

VISION_LLM = "meta-llama/llama-4-scout-17b-16e-instruct"

class ImageDeduplicator:
    """Prevents duplicate images using content hashing"""
    def __init__(self):
        self.image_hashes = set()
        self.hash_to_b64 = {}
    
    def get_image_hash(self, img_path):
        with open(img_path, 'rb') as f:
            return hashlib.sha256(f.read()).hexdigest()
    
    def is_duplicate(self, img_path):
        h = self.get_image_hash(img_path)
        if h in self.image_hashes:
            return True
        self.image_hashes.add(h)
        return False

    def store_b64(self, img_path, b64_data):
        h = self.get_image_hash(img_path)
        self.hash_to_b64[h] = b64_data
        return h

def extract_images_from_pptx(pptx_path, deduplicator=None, output_dir="extracted_data"):
    """Extracts images blindly; deduplication happens later."""
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    img_count = 0
    
    for slide_in, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                img_path = os.path.join(output_dir, f"slide_{slide_in}_img_{img_count}.{image.ext}")
                with open(img_path, "wb") as f:
                    f.write(image.blob)
                img_count += 1
                
    print(f"✓ Extracted {img_count} raw images from PPTX")
    return img_count

def extract_images_from_docx(docx_path, deduplicator=None, output_dir="extracted_data"):
    """Extracts images blindly from the DOCX zip archive."""
    os.makedirs(output_dir, exist_ok=True)
    img_count = 0

    with zipfile.ZipFile(docx_path, 'r') as docx_zip:
        for item in docx_zip.namelist():
            if item.startswith('word/media/') and item.endswith(('.png', '.jpg', '.jpeg')):
                image_data = docx_zip.read(item)
                ext = item.split('.')[-1]
                img_path = os.path.join(output_dir, f"docx_img_{uuid.uuid4().hex[:8]}.{ext}")
                
                with open(img_path, "wb") as f:
                    f.write(image_data)
                img_count += 1

    print(f"✓ Extracted {img_count} raw images from DOCX")
    return img_count

def extract_elements(file_path, deduplicator=None):
    """Extracts text, tables, and images from PDF, DOCX, PPTX, or TXT."""
    print(f"Parsing document: {file_path}...")
    os.makedirs("extracted_data", exist_ok=True)

    raw_elements = partition(
        filename=file_path,
        strategy="hi_res",
        extract_images_in_pdf=True,
        extract_image_block_types=["Image", "Table"],
        extract_image_block_to_payload=False,
        extract_image_block_output_dir="extracted_data",
    )

    if file_path.lower().endswith(".pptx"):
       extract_images_from_pptx(file_path, deduplicator, "extracted_data")
    elif file_path.lower().endswith(".docx"):
        extract_images_from_docx(file_path, deduplicator, "extracted_data")
        
    return raw_elements

def categorize_and_chunk(raw_elements):
    """Categorizes elements and chunks long text to prevent retriever fragmentation."""
    raw_texts = []
    raw_tables = []

    for element in raw_elements:
        if "Table" in str(type(element)):
            raw_tables.append(str(element))
        elif "Text" in str(type(element)) or "ListItem" in str(type(element)) or "Title" in str(type(element)):
            raw_texts.append(str(element))

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunked_texts = text_splitter.split_text(" ".join(raw_texts)) if raw_texts else []

    return chunked_texts, raw_tables

def encode_and_resize_image(img_path, max_size=1024):
    """Resizes images while maintaining aspect ratio, then encodes to Base64."""
    img = Image.open(img_path)
    img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
    buffered = io.BytesIO()
    img.save(buffered, format=img.format or 'PNG')
    return base64.b64encode(buffered.getvalue()).decode("utf-8")

def process_extracted_images(deduplicator=None, directory="extracted_data"):
    """Encodes all extracted images in the directory with deduplication."""
    img_base64_list = []
    processed_hashes = set()
    
    if os.path.exists(directory):
        for img_file in sorted(os.listdir(directory)):
            if img_file.endswith((".jpg", ".png", ".jpeg")):
                img_path = os.path.join(directory, img_file)
                
                if deduplicator and deduplicator.is_duplicate(img_path):
                    continue
                
                try:
                    b64 = encode_and_resize_image(img_path)
                    b64_hash = hashlib.sha256(b64.encode()).hexdigest()
                    
                    if b64_hash not in processed_hashes:
                        img_base64_list.append(b64)
                        processed_hashes.add(b64_hash)
                except Exception as e:
                    print(f"  ✗ Error processing {img_file}: {e}")
    
    return img_base64_list

def generate_summaries(texts, tables, images_b64):
    """Generates retrievable summaries for texts, tables, and images."""
    print("Generating summaries for indexing...")
    llm = ChatGroq(model=VISION_LLM, temperature=0, max_tokens=250)
    
    text_prompt = ChatPromptTemplate.from_template(
        "Give a concise, highly searchable summary of the following content optimized for retrieval: {element}"
    )    
    text_summarizer = {"element": lambda x: x} | text_prompt | llm | StrOutputParser()

    text_summaries = text_summarizer.batch(texts, {"max_concurrency": 5}) if texts else []
    table_summaries = text_summarizer.batch(tables, {"max_concurrency": 5}) if tables else []

    image_summaries = []
    if images_b64:  
        for b64_img in images_b64:
            msg = llm.invoke([
                HumanMessage(content=[
                    {"type": "text", "text": "Describe this image or chart in high detail. Focus on data points, labels, and core concepts so it can be easily searched."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64_img}"}}
                ])
            ])
            image_summaries.append(msg.content)

    return text_summaries, table_summaries, image_summaries